"""
GitHub Issue Triager using Gemini Key Rotation.

Analyzes a GitHub issue and produces severity label, category,
reproduction steps summary, and suggested assignee type.
"""

import argparse
import base64
import json
import os
import re
import requests

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

TRIAGE_PROMPT = """You are a GitHub issue triager. Analyze the issue and return a JSON object with:
{
  "severity": "critical|high|medium|low",
  "category": "bug|feature|documentation|question|performance|security",
  "priority_score": 1-10,
  "labels": ["list", "of", "suggested", "labels"],
  "summary": "one sentence summary",
  "reproduction_clear": true/false,
  "assignee_type": "frontend|backend|devops|documentation|security|any",
  "needs_more_info": true/false,
  "triage_notes": "2-3 sentences of triager notes"
}
Return only valid JSON, no markdown formatting."""


def parse_json_response(text: str) -> dict:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    return json.loads(cleaned)


def triage_issue(title: str, body: str, labels: list[str] = None) -> dict:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)

    issue_text = f"Title: {title}\n\nBody:\n{body}"
    if labels:
        issue_text += f"\n\nExisting labels: {', '.join(labels)}"

    messages = [
        SystemMessage(content=TRIAGE_PROMPT),
        HumanMessage(content=issue_text),
    ]

    response = llm.invoke(messages)
    return parse_json_response(response.content)


def fetch_github_issue(url: str) -> tuple[str, str, list]:
    match = re.match(r"[https://github.com/](https://github.com/)([^/]+)/([^/]+)/issues/(\d+)", url)
    if not match:
        raise ValueError(f"Invalid GitHub issue URL: {url}")

    owner, repo, issue_num = match.groups()
    api_url = f"[https://api.github.com/repos/](https://api.github.com/repos/){owner}/{repo}/issues/{issue_num}"
    headers = {}
    if token := os.getenv("GITHUB_TOKEN"):
        headers["Authorization"] = f"token {token}"

    r = requests.get(api_url, headers=headers, timeout=10)
    r.raise_for_status()
    data = r.json()
    return data["title"], data.get("body", ""), [l["name"] for l in data.get("labels", [])]


def format_triage_markdown(title: str, result: dict) -> str:
    severity = result.get("severity", "medium")
    severity_emoji = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "🟢"}.get(severity, "⚪")
    labels = result.get("labels", [])

    return f"""# GitHub Issue Triage Report: {title}

* **{severity_emoji} Severity:** {severity.upper()} (Priority Score: {result.get('priority_score', 'N/A')}/10)
* **📁 Category:** {result.get('category', 'N/A')}
* **👤 Suggested Assignee:** {result.get('assignee_type', 'any')} team
* **🏷️ Labels:** {', '.join(labels) if labels else 'None'}
* **❓ Needs More Info:** {'Yes' if result.get('needs_more_info') else 'No'}
* **🔍 Reproduction Steps Clear:** {'Yes' if result.get('reproduction_clear') else 'No'}

### Summary
{result.get('summary', 'N/A')}

### Triager Notes
{result.get('triage_notes', 'N/A')}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "triage_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "triage_report.docx") -> str:
    doc = Document()
    doc.add_heading("GitHub Issue Triage Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "triage_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Issue Triage Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="GitHub Issue Triager")
    group = parser.add_mutually_exclusive_group(required=False)
    group.add_argument("--issue-url", help="GitHub issue URL")
    group.add_argument("--title", help="Issue title (use with --body)")
    parser.add_argument("--query", help="General execution query/prompt")
    parser.add_argument("--body", default="", help="Issue body text")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    if args.issue_url:
        title, body, labels = fetch_github_issue(args.issue_url)
    elif args.title:
        title, body, labels = args.title, args.body, []
    elif "github.com" in task_prompt and "/issues/" in task_prompt:
        url_match = re.search(r"https://github\.com/[^\s]+", task_prompt)
        if url_match:
            title, body, labels = fetch_github_issue(url_match.group(0))
        else:
            title, body, labels = "Generic Task Request", task_prompt, []
    else:
        title = "Issue Report Request"
        body = task_prompt or "Login fails on mobile Safari when clicking submit button."
        labels = []

    result = triage_issue(title, body, labels)
    report_md = format_triage_markdown(title, result)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or title, report_md)


if __name__ == "__main__":
    main()
