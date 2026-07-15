"""
Social Media Content Agent using CrewAI and Gemini Key Rotation.

Generates platform-optimized content (Twitter/X, LinkedIn, Instagram)
from a topic or article URL.
"""

import argparse
import base64
import os
import re

from crewai import Agent, Crew, Process, Task
from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


def generate_social_content(topic: str, brand: str, platforms: list[str]) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.7)

    strategist = Agent(
        role="Social Media Strategist",
        goal="Analyze the topic and define the key message, target audience, and tone for each platform",
        backstory="Award-winning social media strategist who has grown 50+ brand accounts to 100k+ followers.",
        llm=llm,
        verbose=False,
    )

    writer = Agent(
        role="Social Media Copywriter",
        goal="Write engaging, platform-optimized content that drives engagement",
        backstory="Viral content creator with expertise in platform-specific formats, hashtags, and hooks.",
        llm=llm,
        verbose=False,
    )

    strategy_task = Task(
        description=f"""Analyze this topic for social media: "{topic}"
Brand: {brand or 'Not specified'}
Platforms: {', '.join(platforms)}
Define: core message, target audience, emotional hook, 5 relevant hashtags.""",
        agent=strategist,
        expected_output="Content strategy: message, audience, hook, and hashtags",
    )

    writing_task = Task(
        description=f"""Write social media posts for: {', '.join(platforms)}
Topic: {topic}. Brand: {brand or 'General'}.

For each platform:
- Twitter/X: 2 tweet variations (under 280 chars each) + thread opener
- LinkedIn: Professional post (150-200 words) with storytelling hook
- Instagram: Caption (100-150 words) + 15 hashtags

Make them platform-native — Twitter punchy, LinkedIn thoughtful, Instagram visual.""",
        agent=writer,
        expected_output="Platform-optimized posts for all requested platforms",
        context=[strategy_task],
    )

    crew = Crew(
        agents=[strategist, writer],
        tasks=[strategy_task, writing_task],
        process=Process.sequential,
        verbose=False,
    )

    return str(crew.kickoff())


def format_report_markdown(topic: str, brand: str, platforms: list[str], content_text: str) -> str:
    brand_display = brand if brand else "General"
    return f"""# Social Media Content Strategy Report

* **Topic:** {topic}
* **Brand:** {brand_display}
* **Target Platforms:** {', '.join(platforms)}

---

{content_text}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "social_media_content_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "social_media_content_report.docx") -> str:
    doc = Document()
    doc.add_heading("Social Media Content Strategy Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "social_media_content_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Social Media Content"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Social Media Content Agent")
    parser.add_argument("--topic", help="Content topic")
    parser.add_argument("--brand", default="", help="Brand name (optional)")
    parser.add_argument("--platforms", help="Comma-separated platforms")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    topic = args.topic or "How AI is transforming software development in 2026"
    platforms_raw = args.platforms or "twitter,linkedin,instagram"
    brand = args.brand

    if task_prompt and not args.topic:
        topic = task_prompt

    platforms = [p.strip() for p in platforms_raw.split(",")]

    content = generate_social_content(topic, brand, platforms)
    report_md = format_report_markdown(topic, brand, platforms, content)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or topic, report_md)


if __name__ == "__main__":
    main()
