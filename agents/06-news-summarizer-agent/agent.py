"""
News Summarizer Agent using Gemini Key Rotation.

Fetches news articles and produces structured summaries with key insights.
"""

import argparse
import base64
import os
import re
import requests

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

NEWS_API_KEY = os.getenv("NEWS_API_KEY")


def fetch_news(topic: str, count: int = 5) -> list[dict]:
    if not NEWS_API_KEY:
        # Return mock data if no API key
        return [
            {"title": f"Major development in {topic}", "description": f"Researchers announce breakthrough in {topic} field.", "url": "https://example.com/1", "source": {"name": "Tech News"}},
            {"title": f"{topic.title()} industry sees rapid growth", "description": f"New report shows {topic} adoption up 40% year-over-year.", "url": "https://example.com/2", "source": {"name": "Business Daily"}},
            {"title": f"Experts weigh in on {topic} challenges", "description": f"Leading experts discuss obstacles facing the {topic} space.", "url": "https://example.com/3", "source": {"name": "Science Weekly"}},
        ]

    url = f"https://newsapi.org/v2/everything?q={topic}&language=en&pageSize={count}&sortBy=publishedAt&apiKey={NEWS_API_KEY}"
    try:
        response = requests.get(url, timeout=10)
        data = response.json()
        return data.get("articles", [])
    except Exception:
        return []


def summarize_news(topic: str, articles: list[dict]) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)

    articles_text = "\n\n".join(
        f"Title: {a.get('title', 'N/A')}\nSource: {a.get('source', {}).get('name', 'Unknown')}\nSummary: {a.get('description', 'N/A')}"
        for a in articles[:5]
    )

    messages = [
        SystemMessage(content=(
            "You are a news analyst. Create a structured news briefing with: "
            "1) Top Story, 2) Key Themes (3 bullet points), 3) What to Watch, 4) Quick Headlines list. "
            "IMPORTANT: Output ONLY the news briefing formatted as Markdown text. Do NOT write Python scripts or meta-explanations."
        )),
        HumanMessage(content=f"Topic: {topic}\n\nArticles:\n{articles_text}"),
    ]

    response = llm.invoke(messages)
    return response.content


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "news_briefing.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "news_briefing.docx") -> str:
    doc = Document()
    doc.add_heading("News Briefing Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "news_briefing.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "News Briefing Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="News Summarizer Agent")
    parser.add_argument("--topic", help="News topic to search")
    parser.add_argument("--query", help="General execution query/prompt")
    parser.add_argument("--count", type=int, default=5, help="Number of articles to fetch")
    args = parser.parse_args()

    topic = args.topic or args.query or os.getenv("TASK_PROMPT") or "artificial intelligence"

    articles = fetch_news(topic, args.count)
    summary = summarize_news(topic, articles)

    print("\n---REPORT_START---")
    print(summary)
    print("---REPORT_END---\n")

    handle_on_demand_export(topic, summary)


if __name__ == "__main__":
    main()
