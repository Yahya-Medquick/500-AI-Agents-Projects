"""
Stock Research Agent using Gemini Key Rotation and Yahoo Finance.

Provides comprehensive stock analysis: price data, financials,
analyst ratings, and AI-powered investment summary.
"""

import argparse
import base64
import os
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

try:
    import yfinance as yf
    HAS_YFINANCE = True
except ImportError:
    HAS_YFINANCE = False


def get_stock_data(ticker: str) -> dict:
    if not HAS_YFINANCE:
        return {"ticker": ticker, "error": "yfinance not installed", "mock": True}

    try:
        stock = yf.Ticker(ticker)
        info = stock.info
    except Exception:
        info = {}

    return {
        "ticker": ticker,
        "name": info.get("longName", ticker),
        "sector": info.get("sector", "N/A"),
        "industry": info.get("industry", "N/A"),
        "price": info.get("currentPrice", info.get("regularMarketPrice", 0)),
        "market_cap": info.get("marketCap", 0),
        "pe_ratio": info.get("trailingPE", "N/A"),
        "forward_pe": info.get("forwardPE", "N/A"),
        "peg_ratio": info.get("pegRatio", "N/A"),
        "revenue_growth": info.get("revenueGrowth", "N/A"),
        "profit_margin": info.get("profitMargins", "N/A"),
        "dividend_yield": info.get("dividendYield", 0),
        "52w_high": info.get("fiftyTwoWeekHigh", "N/A"),
        "52w_low": info.get("fiftyTwoWeekLow", "N/A"),
        "analyst_rating": info.get("recommendationKey", "N/A"),
        "target_price": info.get("targetMeanPrice", "N/A"),
        "description": info.get("longBusinessSummary", "")[:500],
    }


def analyze_stock(data: dict) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)

    stock_info = "\n".join(f"{k}: {v}" for k, v in data.items() if k != "description")

    messages = [
        SystemMessage(content=(
            "You are a financial analyst. Provide a concise stock analysis covering: "
            "Investment Thesis (2-3 sentences), Key Strengths (3 bullets), Key Risks (3 bullets), "
            "Valuation Assessment, and a Verdict (Buy/Hold/Sell with brief reasoning). Keep it under 300 words."
        )),
        HumanMessage(content=f"Analyze this stock:\n{stock_info}\n\nCompany description: {data.get('description', 'N/A')}"),
    ]

    response = llm.invoke(messages)
    return response.content


def format_number(n) -> str:
    if isinstance(n, (int, float)):
        if n >= 1e12:
            return f"${n/1e12:.2f}T"
        if n >= 1e9:
            return f"${n/1e9:.2f}B"
        if n >= 1e6:
            return f"${n/1e6:.2f}M"
        return f"${n:.2f}"
    return str(n)


def format_report_markdown(data: dict, analysis: str) -> str:
    ticker = data.get("ticker", "UNKNOWN")
    name = data.get("name", ticker)
    analyst_rating = str(data.get("analyst_rating", "N/A")).upper()

    return f"""# Stock Research Report: {name} ({ticker})

* **Price:** ${data.get('price', 'N/A')} | **Market Cap:** {format_number(data.get('market_cap', 0))}
* **Sector:** {data.get('sector')} | **Industry:** {data.get('industry')}
* **P/E:** {data.get('pe_ratio')} | **Forward P/E:** {data.get('forward_pe')} | **PEG:** {data.get('peg_ratio')}
* **52W Range:** ${data.get('52w_low')} - ${data.get('52w_high')}
* **Analyst Rating:** {analyst_rating} | **Target Price:** ${data.get('target_price', 'N/A')}

---

## Investment Analysis
{analysis}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "stock_research_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "stock_research_report.docx") -> str:
    doc = Document()
    doc.add_heading("Stock Research Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "stock_research_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Stock Analysis Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Stock Research Agent")
    parser.add_argument("--ticker", help="Stock ticker symbol (e.g., AAPL)")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""
    
    ticker = args.ticker
    if not ticker:
        # Extract ticker symbol pattern (e.g. AAPL, NVDA) from prompt or default
        match = re.search(r"\b[A-Z]{1,5}\b", task_prompt)
        ticker = match.group(0) if match else "AAPL"

    data = get_stock_data(ticker)
    analysis = analyze_stock(data)
    report_md = format_report_markdown(data, analysis)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or ticker, report_md)


if __name__ == "__main__":
    main()
