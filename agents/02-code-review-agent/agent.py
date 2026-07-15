"""
Code Review Agent using LangChain and Gemini Key Rotation.

Reviews code for bugs, security issues, style violations, and
suggests improvements. Accepts a file path, inline code snippet, or query string.
"""

import argparse
import base64
import os
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the rotator utility at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

SYSTEM_PROMPT = """You are an expert code reviewer. Analyze the provided code or task request and return a structured review covering:

1. **Bugs & Correctness** — logic errors, edge cases, exception handling
2. **Security Issues** — injection risks, secrets exposure, unsafe operations
3. **Performance** — inefficiencies, unnecessary computation, memory issues
4. **Code Style** — PEP 8 violations, naming conventions, readability
5. **Improvements** — refactoring suggestions, better patterns

Format: Use markdown. Rate overall quality as: 🟢 Good / 🟡 Needs Work / 🔴 Critical Issues.
IMPORTANT: Output ONLY the review content as formatted Markdown. Do NOT write meta-explanations."""


def review_code(code: str, language: str = "python") -> str:
    # Uses automatic key rotation from GEMINI_API_KEY secret
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    
    messages = [
        SystemMessage(content=SYSTEM_PROMPT),
        HumanMessage(content=f"Review this {language} code or code request:\n\n```{language}\n{code}\n```"),
    ]
    response = llm.invoke(messages)
    return response.content


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "code_review.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "code_review.docx") -> str:
    doc = Document()
    doc.add_heading("Code Review Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "code_review.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Code Review Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Code Review Agent")
    group = parser.add_mutually_exclusive_group(required=False)
    group.add_argument("--file", help="Path to file to review")
    group.add_argument("--code", help="Inline code snippet to review")
    parser.add_argument("--query", help="General review query/prompt")
    parser.add_argument("--language", default="python", help="Programming language (default: python)")
    args = parser.parse_args()

    raw_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    if args.file:
        with open(args.file) as f:
            code = f.read()
    elif args.code:
        code = args.code
    elif raw_prompt:
        code = raw_prompt
    else:
        code = "def add(a, b): return a + b"

    review = review_code(code, args.language)

    print("\n---REPORT_START---")
    print(review)
    print("---REPORT_END---\n")

    handle_on_demand_export(raw_prompt, review)


if __name__ == "__main__":
    main()
