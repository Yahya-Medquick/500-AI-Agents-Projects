"""
Data Analysis Agent using LangChain, pandas, and Gemini Key Rotation.

Loads a CSV/Excel file and answers analytical questions about it using
natural language by executing generated pandas operations.
"""

import argparse
import base64
import os
import random
import re
import pandas as pd

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_experimental.agents import create_pandas_dataframe_agent
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


def create_sample_data(path: str) -> pd.DataFrame:
    """Creates a sample sales dataset for demo."""
    from datetime import date, timedelta

    random.seed(42)
    rows = []
    products = ["Laptop", "Phone", "Tablet", "Monitor", "Keyboard"]
    regions = ["North", "South", "East", "West"]
    start = date(2024, 1, 1)

    for i in range(200):
        d = start + timedelta(days=random.randint(0, 364))
        rows.append({
            "date": d.isoformat(),
            "product": random.choice(products),
            "region": random.choice(regions),
            "quantity": random.randint(1, 20),
            "unit_price": round(random.uniform(50, 2000), 2),
            "revenue": 0,
        })

    df = pd.DataFrame(rows)
    df["revenue"] = df["quantity"] * df["unit_price"]
    df.to_csv(path, index=False)
    return df


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "data_analysis_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "data_analysis_report.docx") -> str:
    doc = Document()
    doc.add_heading("Data Analysis Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "data_analysis_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Analysis Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Data Analysis Agent")
    parser.add_argument("--file", default="sample_data.csv", help="CSV or Excel file to analyze")
    parser.add_argument("--question", help="Analytical question to ask")
    parser.add_argument("--query", help="General execution task/prompt")
    args = parser.parse_args()

    question = args.question or args.query or os.getenv("TASK_PROMPT") or "Summarize the key revenue insights from this dataset."

    if args.file == "sample_data.csv" and not os.path.exists("sample_data.csv"):
        df = create_sample_data("sample_data.csv")
    else:
        ext = os.path.splitext(args.file)[1].lower()
        df = pd.read_excel(args.file) if ext in (".xlsx", ".xls") else pd.read_csv(args.file)

    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    
    agent = create_pandas_dataframe_agent(
        llm,
        df,
        verbose=False,
        allow_dangerous_code=True,
    )

    result = agent.invoke({"input": question})
    output_text = result.get("output", str(result))

    print("\n---REPORT_START---")
    print(output_text)
    print("---REPORT_END---\n")

    handle_on_demand_export(question, output_text)


if __name__ == "__main__":
    main()
