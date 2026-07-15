"""
Web Research Agent using LangGraph and Gemini Key Rotation.

Multi-step agent that conducts web research:
1. Formulates search queries based on the prompt
2. Executes web searches and gathers information
3. Synthesizes findings into a detailed research report

Usage:
    python agent.py --query "Latest advancements in quantum computing"
"""

import argparse
import base64
import os
import sys
from typing import Annotated, TypedDict

# Fix for ModuleNotFoundError when running scripts directly
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "../..")))

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_community.tools import TavilySearchResults
from langchain_core.messages import HumanMessage, SystemMessage
from langgraph.graph import END, StateGraph
from langgraph.graph.message import add_messages
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


class ResearchState(TypedDict):
    messages: Annotated[list, add_messages]
    query: str
    search_queries: list[str]
    search_results: list[str]
    final_report: str


def generate_queries(state: ResearchState) -> ResearchState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    response = llm.invoke([
        SystemMessage(content="You are a research assistant. Generate 3 distinct search queries to thoroughly research the given topic. Return them as a comma-separated list."),
        HumanMessage(content=f"Topic: {state['query']}"),
    ])
    queries = [q.strip() for q in response.content.split(",")][:3]
    return {"search_queries": queries, "messages": [response]}


def execute_searches(state: ResearchState) -> ResearchState:
    tool = TavilySearchResults(max_results=3)
    results = []
    for query in state["search_queries"]:
        try:
            res = tool.invoke({"query": query})
            results.append(f"Query: {query}\nResults: {res}")
        except Exception as e:
            results.append(f"Query: {query}\nFailed with error: {str(e)}")
    return {"search_results": results}


def synthesize_report(state: ResearchState) -> ResearchState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    results_text = "\n\n".join(state["search_results"])

    response = llm.invoke([
        SystemMessage(content="""You are a senior research analyst. Synthesize the provided search results into a comprehensive research report.
Include:
1. Executive Summary
2. Key Findings & Analysis
3. Implications & Trends
4. Conclusion & Key Takeaways"""),
        HumanMessage(content=f"Research Topic: {state['query']}\n\nSearch Information:\n{results_text}"),
    ])

    return {"final_report": response.content, "messages": [response]}


def build_graph():
    graph = StateGraph(ResearchState)
    graph.add_node("plan", generate_queries)
    graph.add_node("search", execute_searches)
    graph.add_node("report", synthesize_report)
    graph.set_entry_point("plan")
    graph.add_edge("plan", "search")
    graph.add_edge("search", "report")
    graph.add_edge("report", END)
    return graph.compile()


def format_report_markdown(query: str, search_queries: list[str], report_text: str) -> str:
    return f"""# Web Research Report

* **Research Topic:** {query}
* **Search Queries Used:** {', '.join(search_queries)}

---

{report_text}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "web_research_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "web_research_report.docx") -> str:
    doc = Document()
    doc.add_heading("Web Research Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "web_research_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Web Research Report"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Web Research Agent")
    parser.add_argument("--query", help="Research query or topic")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or "Latest AI developments"

    agent = build_graph()
    result = agent.invoke({
        "query": task_prompt,
        "messages": [],
        "search_queries": [],
        "search_results": [],
        "final_report": "",
    })

    report_md = format_report_markdown(task_prompt, result["search_queries"], result["final_report"])

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt, report_md)


if __name__ == "__main__":
    main()
