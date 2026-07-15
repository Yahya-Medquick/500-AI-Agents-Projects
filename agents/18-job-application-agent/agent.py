"""
Job Application Agent using CrewAI and Gemini Key Rotation.

Analyzes a job description and a candidate profile, then generates:
- Tailored cover letter
- Resume bullet points to highlight
- Interview preparation questions

Usage:
    python agent.py --job-desc "Senior Python Engineer at Stripe..." --candidate "7 years Python, FastAPI..."
"""

import argparse
import base64
import os
import re

from crewai import Agent, Crew, Process, Task
from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

SAMPLE_JOB = """Senior Python Engineer at Stripe
We're looking for a Senior Python Engineer to join our API Platform team.

Requirements:
- 5+ years Python development
- Experience with distributed systems
- Strong understanding of REST APIs and microservices
- Experience with PostgreSQL, Redis
- Kubernetes experience preferred
- Strong communication skills

Responsibilities:
- Design and build high-performance APIs handling millions of requests/day
- Lead technical design reviews
- Mentor junior engineers
- Collaborate with product managers on technical feasibility
"""

SAMPLE_CANDIDATE = """
Jane Doe — 7 years Python experience
Current role: Senior Software Engineer at DataCorp
Skills: Python, FastAPI, Django, PostgreSQL, Redis, Docker, Kubernetes, AWS
Achievements:
- Built API platform handling 5M requests/day
- Led team of 4 engineers
- Reduced API latency by 40%
- Mentored 3 junior engineers
Education: BS Computer Science, UC Berkeley
"""


def run_job_application_crew(job_desc: str, candidate_profile: str) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.4)

    analyst = Agent(
        role="Job Requirements Analyst",
        goal="Analyze the job description and identify key requirements, values, and culture signals",
        backstory="Ex-hiring manager at FAANG with 10 years recruiting experience. Expert at decoding job descriptions.",
        llm=llm,
        verbose=False,
    )

    writer = Agent(
        role="Career Coach and Application Writer",
        goal="Create tailored application materials that maximize interview chances",
        backstory="Career coach who has helped 500+ candidates land roles at top tech companies.",
        llm=llm,
        verbose=False,
    )

    analyst_task = Task(
        description=f"""Analyze this job description:
{job_desc}

Extract: top 5 required skills, culture signals, what this company values most, potential red flags, and key phrases to mirror in the application.""",
        agent=analyst,
        expected_output="Job analysis: key requirements, culture signals, important keywords",
    )

    application_task = Task(
        description=f"""Using the job analysis, create application materials for this candidate:
{candidate_profile}

Produce:
1. COVER LETTER (250-300 words, 3 paragraphs: hook, evidence, close)
2. TOP 5 RESUME BULLETS TO HIGHLIGHT (tailored to this specific role)
3. 10 LIKELY INTERVIEW QUESTIONS (5 behavioral, 5 technical) with suggested answer frameworks
4. NEGOTIATION RANGE ESTIMATE based on role seniority and company""",
        agent=writer,
        expected_output="Cover letter, resume bullets, interview questions, salary range",
        context=[analyst_task],
    )

    crew = Crew(
        agents=[analyst, writer],
        tasks=[analyst_task, application_task],
        process=Process.sequential,
        verbose=False,
    )

    return str(crew.kickoff())


def format_report_markdown(job_desc: str, candidate_profile: str, application_materials: str) -> str:
    return f"""# Job Application Package Report

### Application Package Details
{application_materials}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "job_application_package.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "job_application_package.docx") -> str:
    doc = Document()
    doc.add_heading("Job Application Package Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "job_application_package.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Job Application Package"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Job Application Agent")
    parser.add_argument("--job-desc", help="Job description text")
    parser.add_argument("--candidate", help="Candidate profile summary")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    job_desc = args.job_desc or SAMPLE_JOB
    candidate_profile = args.candidate or SAMPLE_CANDIDATE

    if task_prompt and not args.job_desc:
        job_desc = task_prompt

    result = run_job_application_crew(job_desc, candidate_profile)
    report_md = format_report_markdown(job_desc, candidate_profile, result)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or "job_application", report_md)


if __name__ == "__main__":
    main()
