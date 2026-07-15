"""
Multi-Agent Debate System using AutoGen-style orchestration and Gemini Key Rotation.

Two AI agents debate a topic from opposing sides, moderated by a judge
who declares a winner and synthesizes the key arguments.

Usage:
    python agent.py --topic "AI will eliminate more jobs than it creates"
    python agent.py --topic "Remote work is better than office work" --rounds 3
"""

import argparse
import base64
import os
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


class DebateAgent:
    def __init__(self, name: str, position: str, expertise: str):
        self.name = name
        self.position = position
        self.expertise = expertise
        # Uses key rotator with gemini-3.5-flash
        self.llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.6)
        self.arguments = []

    def make_argument(self, topic: str, round_num: int, opponent_last_arg: str = "") -> str:
        system_msg = f"""You are {self.name}, a {self.expertise}.
You are arguing {self.position} on this topic.
Make compelling, evidence-based arguments. Be direct and persuasive.
Keep response under 150 words. Round {round_num}."""

        if opponent_last_arg:
            user_msg = f"Topic: {topic}\n\nYour opponent just said: '{opponent_last_arg}'\n\nRespond and advance your argument:"
        else:
            user_msg = f"Topic: {topic}\n\nMake your opening argument for {self.position.upper()}:"

        response = self.llm.invoke([
            SystemMessage(content=system_msg),
            HumanMessage(content=user_msg),
        ])
        argument = response.content
        self.arguments.append(argument)
        return argument


class DebateJudge:
    def __init__(self):
        # Uses key rotator with gemini-3.5-flash
        self.llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)

    def evaluate(self, topic: str, pro_agent: DebateAgent, con_agent: DebateAgent) -> dict:
        pro_args = "\n\n".join(f"Round {i+1}: {a}" for i, a in enumerate(pro_agent.arguments))
        con_args = "\n\n".join(f"Round {i+1}: {a}" for i, a in enumerate(con_agent.arguments))

        response = self.llm.invoke([
            SystemMessage(content="""You are an impartial debate judge. Evaluate both sides fairly.
Return a structured verdict with: winner, score (out of 10 each), strongest argument per side, key insights, and balanced synthesis conclusion."""),
            HumanMessage(content=f"""Topic: "{topic}"

PRO arguments ({pro_agent.name}):
{pro_args}

CON arguments ({con_agent.name}):
{con_args}

Provide your verdict:"""),
        ])
        return {"verdict": response.content}


def format_report_markdown(topic: str, pro: DebateAgent, con: DebateAgent, verdict: str) -> str:
    md_content = f"# Multi-Agent Debate Report\n\n"
    md_content += f"**Topic:** {topic}\n\n"
    md_content += f"* **FOR:** {pro.name} ({pro.expertise})\n"
    md_content += f"* **AGAINST:** {con.name} ({con.expertise})\n\n---\n\n"

    md_content += "## 💬 Debate Transcripts\n\n"
    num_rounds = max(len(pro.arguments), len(con.arguments))
    for i in range(num_rounds):
        md_content += f"### Round {i+1}\n\n"
        if i < len(pro.arguments):
            md_content += f"🟢 **{pro.name} (FOR):**\n{pro.arguments[i]}\n\n"
        if i < len(con.arguments):
            md_content += f"🔴 **{con.name} (AGAINST):**\n{con.arguments[i]}\n\n"

    md_content += "---\n\n## 🏛️ Judge's Verdict\n\n"
    md_content += verdict + "\n"
    return md_content


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "multi_agent_debate_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "multi_agent_debate_report.docx") -> str:
    doc = Document()
    doc.add_heading("Multi-Agent Debate Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "multi_agent_debate_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Multi-Agent Debate Report"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def run_debate(topic: str, rounds: int = 2) -> tuple[DebateAgent, DebateAgent, str]:
    pro = DebateAgent(
        name="Dr. Alex Chen",
        position="FOR",
        expertise="technology economist and AI researcher"
    )
    con = DebateAgent(
        name="Prof. Sarah Martinez",
        position="AGAINST",
        expertise="labor economist and social policy expert"
    )
    judge = DebateJudge()

    last_con_arg = ""

    for round_num in range(1, rounds + 1):
        pro_arg = pro.make_argument(topic, round_num, last_con_arg)
        con_arg = con.make_argument(topic, round_num, pro_arg)
        last_con_arg = con_arg

    verdict_data = judge.evaluate(topic, pro, con)
    return pro, con, verdict_data.get("verdict", "")


def main():
    parser = argparse.ArgumentParser(description="Multi-Agent Debate System")
    parser.add_argument("--topic", help="Debate topic")
    parser.add_argument("--rounds", type=int, default=2, help="Number of debate rounds (1-4)")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    topic = args.topic or (task_prompt if task_prompt else "AI will create more jobs than it eliminates over the next decade")
    rounds = max(1, min(4, args.rounds))

    pro, con, verdict = run_debate(topic, rounds)
    report_md = format_report_markdown(topic, pro, con, verdict)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or topic, report_md)


if __name__ == "__main__":
    main()
