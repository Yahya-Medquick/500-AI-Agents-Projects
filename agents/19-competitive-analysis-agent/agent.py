"""
Competitive Analysis Agent using LangGraph and Gemini Key Rotation.

Multi-step agent that analyzes competitors:
1. Identifies key competitors
2. Analyzes each competitor's strengths/weaknesses
3. Generates competitive positioning recommendations

Usage:
    python agent.py --company "Notion" --industry "productivity software"
"""

import argparse
import base64
import os
import re
from typing import Annotated, TypedDict

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langgraph.graph import END, StateGraph
from langgraph.graph.message import add_messages
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


class AnalysisState(TypedDict):
    messages: Annotated[list, add_messages]
    company: str
    industry: str
    competitors: list[str]
    competitor_analyses: dict[str, str]
    final_report: str


def identify_competitors(state: AnalysisState) -> AnalysisState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    response = llm.invoke([
        SystemMessage(content="You are a market research analyst. List exactly 5 main competitors as a comma-separated list. Nothing else."),
        HumanMessage(content=f"Company: {state['company']}\nIndustry: {state['industry']}\n\nList 5 main competitors:"),
    ])
    competitors = [c.strip() for c in response.content.split(",")][:5]
    return {"competitors": competitors, "messages": [response]}


def analyze_competitor(state: AnalysisState) -> AnalysisState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    analyses = {}

    for competitor in state["competitors"]:
        response = llm.invoke([
            SystemMessage(content="Provide a concise competitive analysis in 100 words covering: main products, strengths (2), weaknesses (2), pricing model, target market."),
            HumanMessage(content=f"Analyze {competitor} vs {state['company']} in {state['industry']}:"),
        ])
        analyses[competitor] = response.content

    return {"competitor_analyses": analyses}


def generate_report(state: AnalysisState) -> AnalysisState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)

    analyses_text = "\n\n".join(
        f"**{name}:**\n{analysis}"
        for name, analysis in state["competitor_analyses"].items()
    )

    response = llm.invoke([
        SystemMessage(content="""You are a strategic consultant. Create a competitive analysis report with:
1. Executive Summary (3 sentences)
2. Competitive Landscape Table (company, strength, weakness, price)
3. Market Gaps & Opportunities (3 bullet points)
4. Strategic Recommendations for {company} (5 action items)
5. Threat Assessment (High/Medium/Low for each competitor)""".replace("{company}", state["company"])),
        HumanMessage(content=f"Company: {state['company']}\nIndustry: {state['industry']}\n\nCompetitor analyses:\n{analyses_text}"),
    ])

    return {"final_report": response.content, "messages": [response]}


def build_graph():
    graph = StateGraph(AnalysisState)
    graph.add_node("identify", identify_competitors)
    graph.add_node("analyze", analyze_competitor)
    graph.add_node("report", generate_report)
    graph.set_entry_point("identify")
    graph.add_edge("identify", "analyze")
    graph.add_edge("analyze", "report")
    graph.add_edge("report", END)
    return graph.compile()


def format_report_markdown(company: str, industry: str, competitors: list[str], report_text: str) -> str:
    return f"""# Competitive Analysis Report: {company}

* **Target Company:** {company}
* **Industry:** {industry}
* **Identified Competitors:** {', '.join(competitors)}

---

{report_text}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "competitive_analysis_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "competitive_analysis_report.docx") -> str:
    doc = Document()
    doc.add_heading("Competitive Analysis Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "competitive_analysis_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Competitive Analysis"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Competitive Analysis Agent")
    parser.add_argument("--company", help="Company to analyze")
    parser.add_argument("--industry", help="Industry")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    company = args.company or "Notion"
    industry = args.industry or "productivity and collaboration software"

    if task_prompt and not args.company:
        company = task_prompt

    agent = build_graph()
    result = agent.invoke({
        "company": company,
        "industry": industry,
        "messages": [],
        "competitors": [],
        "competitor_analyses": {},
        "final_report": "",
    })

    report_md = format_report_markdown(company, industry, result["competitors"], result["final_report"])

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or company, report_md)


if __name__ == "__main__":
    main()
