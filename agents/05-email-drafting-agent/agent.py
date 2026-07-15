"""
Email Drafting Agent using CrewAI and Gemini Key Rotation.

A two-agent crew that drafts professional emails:
- Analyst agent: understands context and tone requirements
- Writer agent: drafts the final email
"""

import argparse
import base64
import os
import re

from crewai import Agent, Crew, Process, Task
from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


def build_email_crew(context: str, tone: str, recipient: str) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.3)

    analyst = Agent(
        role="Email Context Analyst",
        goal="Understand the email context, extract key points, and define the structure",
        backstory="You are an expert business communication analyst who distills complex situations into clear email requirements.",
        llm=llm,
        verbose=False,
    )

    writer = Agent(
        role="Professional Email Writer",
        goal="Draft clear, concise, and effective professional emails",
        backstory="You are a professional copywriter specializing in business emails that get responses.",
        llm=llm,
        verbose=False,
    )

    analyze_task = Task(
        description=f"""Analyze this email requirement:
Context: {context}
Recipient: {recipient}
Desired tone: {tone}

Extract: purpose, key points to cover, call to action, subject line suggestion.""",
        agent=analyst,
        expected_output="Structured email brief: purpose, key points, CTA, and suggested subject line",
    )

    write_task = Task(
        description=f"""Using the analysis, draft a complete professional email.
Tone: {tone}. Recipient: {recipient}.
Include: Subject line, greeting, body paragraphs, closing, signature placeholder.
Keep it concise — under 200 words for the body. Output ONLY formatted Markdown text.""",
        agent=writer,
        expected_output="Complete formatted email ready to send",
        context=[analyze_task],
    )

    crew = Crew(
        agents=[analyst, writer],
        tasks=[analyze_task, write_task],
        process=Process.sequential,
        verbose=False,
    )

    result = crew.kickoff()
    return str(result)


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "email_draft.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "email_draft.docx") -> str:
    doc = Document()
    doc.add_heading("Email Draft", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "email_draft.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Email Draft Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Email Drafting Agent")
    parser.add_argument("--context", help="Email context/purpose")
    parser.add_argument("--query", help="General execution task")
    parser.add_argument("--tone", default="professional and friendly", help="Email tone")
    parser.add_argument("--recipient", default="a potential client", help="Who the email is for")
    args = parser.parse_args()

    context = (
        args.context
        or args.query
        or os.getenv("TASK_PROMPT")
        or "Follow up on our product demo from last Tuesday. They seemed interested but haven't responded."
    )

    email = build_email_crew(context, args.tone, args.recipient)

    print("\n---REPORT_START---")
    print(email)
    print("---REPORT_END---\n")

    handle_on_demand_export(context, email)


if __name__ == "__main__":
    main()
