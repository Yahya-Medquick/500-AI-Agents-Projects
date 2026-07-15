"""
Customer Support Agent using LangGraph with RAG and Gemini Key Rotation.

Handles customer queries using a knowledge base (product docs).
Routes complex issues to human escalation.
"""

import argparse
import base64
import os
import re
from pathlib import Path
from typing import Annotated, Literal, TypedDict

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_community.vectorstores import FAISS
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langgraph.graph import END, StateGraph
from langgraph.graph.message import add_messages
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

SAMPLE_KB = [
    "Product: CloudSync Pro. Features: real-time sync across 5 devices, 1TB storage, offline mode, version history 30 days.",
    "Pricing: Basic $9/mo (100GB, 2 devices), Pro $19/mo (1TB, 5 devices), Business $49/mo (5TB, unlimited devices).",
    "Cancellation: Cancel anytime from Account > Subscription > Cancel. Refunds available within 14 days of charge.",
    "Password reset: Go to login page, click 'Forgot Password', enter email. Reset link expires in 1 hour.",
    "Sync issues: Check internet connection, ensure app is updated, try Sign Out and Sign In. If persists, contact support.",
    "Supported platforms: Windows 10+, macOS 12+, iOS 15+, Android 10+, Linux (Beta).",
    "Data security: AES-256 encryption at rest and in transit. SOC 2 Type II certified. Zero-knowledge architecture.",
    "File size limit: Individual files up to 10GB (Pro/Business), 2GB (Basic). No limit on total number of files.",
]

ESCALATION_KEYWORDS = ["refund", "lawsuit", "furious", "fraud", "broken", "data loss", "cancel account", "charge", "billing error"]


class SupportState(TypedDict):
    messages: Annotated[list, add_messages]
    user_input: str
    retrieved_context: str
    response: str
    escalate: bool


def retrieve_context(state: SupportState) -> SupportState:
    query = state["user_input"]
    if not hasattr(retrieve_context, "vectorstore"):
        texts = getattr(retrieve_context, "kb_texts", SAMPLE_KB)
        splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=20)
        docs_split = splitter.create_documents(texts)
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        retrieve_context.vectorstore = FAISS.from_documents(docs_split, embeddings)

    docs = retrieve_context.vectorstore.similarity_search(query, k=3)
    context = "\n".join(d.page_content for d in docs)
    return {"retrieved_context": context}


def check_escalation(state: SupportState) -> SupportState:
    text = state["user_input"].lower()
    needs_escalation = any(kw in text for kw in ESCALATION_KEYWORDS)
    return {"escalate": needs_escalation}


def generate_response(state: SupportState) -> SupportState:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.2)
    conversation = state["messages"][:-1]  # exclude latest user msg

    if state.get("escalate"):
        response_text = (
            "I understand your concern and I want to make sure this gets the attention it deserves. "
            "I'm connecting you with a senior support specialist who can resolve this directly. "
            f"You'll hear back within 2 hours. Your case ID is #{abs(hash(state['user_input'])) % 100000}."
        )
    else:
        messages = [
            SystemMessage(content=f"""You are a helpful customer support agent for CloudSync Pro.
Use this knowledge base context to answer accurately:
{state['retrieved_context']}

Be friendly, concise, and solution-focused. If unsure, say so honestly."""),
            *conversation,
            HumanMessage(content=state["user_input"]),
        ]
        response = llm.invoke(messages)
        response_text = response.content

    return {"response": response_text, "messages": [AIMessage(content=response_text)]}


def build_graph():
    graph = StateGraph(SupportState)
    graph.add_node("retrieve", retrieve_context)
    graph.add_node("check_escalation", check_escalation)
    graph.add_node("generate", generate_response)
    graph.set_entry_point("retrieve")
    graph.add_edge("retrieve", "check_escalation")
    graph.add_edge("check_escalation", "generate")
    graph.add_edge("generate", END)
    return graph.compile()


def load_kb_texts(kb_dir: str | None) -> list[str]:
    if not kb_dir:
        return SAMPLE_KB

    root = Path(kb_dir)
    if not root.is_dir():
        raise ValueError(f"Knowledge base directory does not exist: {kb_dir}")

    texts = []
    for path in sorted(root.rglob("*")):
        if path.is_file() and path.suffix.lower() in {".txt", ".md"}:
            texts.append(path.read_text(encoding="utf-8"))

    if not texts:
        raise ValueError(f"No .txt or .md files found in knowledge base directory: {kb_dir}")

    return texts


def format_report_markdown(query: str, response: str, escalated: bool) -> str:
    status = "ESCALATED TO HUMAN SUPPORT" if escalated else "RESOLVED VIA KB"
    return f"""# Customer Support Query Report

* **Query:** {query}
* **Status:** {status}

### Response
{response}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "customer_support_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "customer_support_report.docx") -> str:
    doc = Document()
    doc.add_heading("Customer Support Query Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "customer_support_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Customer Support Resolution"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Customer Support Agent")
    parser.add_argument("--kb-dir", help="Directory containing .txt or .md support knowledge base files")
    parser.add_argument("--query", help="Customer query prompt/task")
    args = parser.parse_args()

    retrieve_context.kb_texts = load_kb_texts(args.kb_dir)
    if hasattr(retrieve_context, "vectorstore"):
        delattr(retrieve_context, "vectorstore")

    agent = build_graph()
    task_prompt = args.query or os.getenv("TASK_PROMPT")

    if task_prompt:
        state = {
            "messages": [HumanMessage(content=task_prompt)],
            "user_input": task_prompt,
            "retrieved_context": "",
            "response": "",
            "escalate": False,
        }
        res_state = agent.invoke(state)
        report_md = format_report_markdown(task_prompt, res_state["response"], res_state.get("escalate", False))

        print("\n---REPORT_START---")
        print(report_md)
        print("---REPORT_END---\n")

        handle_on_demand_export(task_prompt, report_md)
    else:
        state = {"messages": [], "user_input": "", "retrieved_context": "", "response": "", "escalate": False}
        print("\n🎧 Customer Support Agent (CloudSync Pro)")
        print("Type 'quit' to exit\n")

        while True:
            user_input = input("Customer: ").strip()
            if user_input.lower() in ("quit", "exit", "q"):
                break
            if not user_input:
                continue

            state["user_input"] = user_input
            state["messages"].append(HumanMessage(content=user_input))
            state = agent.invoke(state)

            escalation_indicator = " [ESCALATED]" if state.get("escalate") else ""
            print(f"\nAgent{escalation_indicator}: {state['response']}\n")


if __name__ == "__main__":
    main()
