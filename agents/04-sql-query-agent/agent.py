"""
SQL Query Agent using LangChain and Gemini Key Rotation.

Connects to a SQLite database and answers natural language questions
by generating and executing SQL queries.
"""

import argparse
import base64
import os
import re
import sqlite3
from urllib.parse import quote

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain.agents import create_sql_agent
from langchain.agents.agent_toolkits import SQLDatabaseToolkit
from langchain.agents.agent_types import AgentType
from langchain_community.utilities import SQLDatabase
from pptx import Presentation

# Import the rotator utility at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


def create_demo_database(db_path: str):
    """Creates a demo e-commerce SQLite database for testing."""
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.executescript("""
        CREATE TABLE IF NOT EXISTS customers (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            email TEXT UNIQUE,
            country TEXT,
            created_at DATE DEFAULT CURRENT_DATE
        );
        CREATE TABLE IF NOT EXISTS products (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            category TEXT,
            price REAL NOT NULL,
            stock INTEGER DEFAULT 0
        );
        CREATE TABLE IF NOT EXISTS orders (
            id INTEGER PRIMARY KEY,
            customer_id INTEGER REFERENCES customers(id),
            product_id INTEGER REFERENCES products(id),
            quantity INTEGER NOT NULL,
            total REAL NOT NULL,
            order_date DATE DEFAULT CURRENT_DATE
        );
        INSERT OR IGNORE INTO customers VALUES
            (1,'Alice Johnson','alice@example.com','USA','2024-01-15'),
            (2,'Bob Smith','bob@example.com','UK','2024-02-20'),
            (3,'Carlos Lima','carlos@example.com','Brazil','2024-03-10'),
            (4,'Diana Prince','diana@example.com','USA','2024-01-05');
        INSERT OR IGNORE INTO products VALUES
            (1,'Laptop Pro','Electronics',1299.99,45),
            (2,'Wireless Mouse','Electronics',29.99,200),
            (3,'Python Book','Books',49.99,120),
            (4,'Standing Desk','Furniture',599.99,15);
        INSERT OR IGNORE INTO orders VALUES
            (1,1,1,1,1299.99,'2024-04-01'),
            (2,1,2,2,59.98,'2024-04-01'),
            (3,2,3,1,49.99,'2024-04-05'),
            (4,3,4,1,599.99,'2024-04-10'),
            (5,4,1,1,1299.99,'2024-04-12'),
            (6,2,2,3,89.97,'2024-04-15');
    """)
    conn.commit()
    conn.close()


def sqlite_uri(db_path: str, read_only: bool = True) -> str:
    abs_path = os.path.abspath(db_path)
    if read_only:
        return f"sqlite:///file:{quote(abs_path)}?mode=ro&uri=true"
    return f"sqlite:///{abs_path}"


def build_agent(db_path: str, read_only: bool = True):
    db = SQLDatabase.from_uri(sqlite_uri(db_path, read_only=read_only))
    
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    
    toolkit = SQLDatabaseToolkit(db=db, llm=llm)
    agent = create_sql_agent(
        llm=llm,
        toolkit=toolkit,
        agent_type=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
        verbose=False,
    )
    return agent, db


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "sql_query_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "sql_query_report.docx") -> str:
    doc = Document()
    doc.add_heading("SQL Analysis Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "sql_query_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "SQL Query Result Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="SQL Query Agent")
    parser.add_argument("--db", default="demo.sqlite", help="SQLite database path")
    parser.add_argument("--question", help="Natural language question")
    parser.add_argument("--query", help="General task query")
    parser.add_argument("--allow-write", action="store_true", help="Open DB read-write")
    args = parser.parse_args()

    question = args.question or args.query or os.getenv("TASK_PROMPT") or "How many total customers are in the database?"

    if args.db == "demo.sqlite" and not os.path.exists("demo.sqlite"):
        create_demo_database("demo.sqlite")

    agent, db = build_agent(args.db, read_only=not args.allow_write)

    result = agent.invoke({"input": question})
    output_text = result.get("output", str(result))

    print("\n---REPORT_START---")
    print(output_text)
    print("---REPORT_END---\n")

    handle_on_demand_export(question, output_text)


if __name__ == "__main__":
    main()
