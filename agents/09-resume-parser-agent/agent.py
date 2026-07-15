"""
Resume Parser Agent using Gemini Key Rotation.

Extracts structured information from resume text or PDF:
contact info, skills, experience, education, and provides
a candidate summary and fit score for a job description.
"""

import argparse
import base64
import json
import os
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

PARSE_PROMPT = """Extract structured information from this resume and return JSON:
{
  "name": "full name",
  "email": "email or null",
  "phone": "phone or null",
  "location": "city, country or null",
  "linkedin": "URL or null",
  "github": "URL or null",
  "summary": "2-3 sentence professional summary",
  "years_experience": number,
  "current_title": "current/most recent job title",
  "skills": {
    "languages": ["Python", "JavaScript"],
    "frameworks": ["Django", "React"],
    "tools": ["Docker", "Git"],
    "soft_skills": ["leadership"]
  },
  "experience": [{"title": "...", "company": "...", "duration": "...", "highlights": ["..."]}],
  "education": [{"degree": "...", "institution": "...", "year": "..."}],
  "certifications": ["..."],
  "languages_spoken": ["English"]
}
Return only valid JSON, no markdown formatting."""

FIT_PROMPT = """Given this candidate profile and job description, return JSON:
{
  "fit_score": 0-100,
  "fit_label": "Excellent|Good|Fair|Poor",
  "strengths": ["matching point 1", "matching point 2"],
  "gaps": ["missing skill 1"],
  "recommendation": "Hire|Consider|Pass",
  "recommendation_reason": "2-3 sentence explanation"
}
Return only valid JSON, no markdown formatting."""


def parse_json_response(text: str) -> dict:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    return json.loads(cleaned)


def read_resume_text(path: str) -> str:
    if path.endswith(".pdf"):
        try:
            import pypdf
            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                return "\n".join(page.extract_text() for page in reader.pages)
        except ImportError:
            raise ImportError("pypdf is required to read PDF files. Install with: pip install pypdf")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_resume(text: str) -> dict:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    messages = [SystemMessage(content=PARSE_PROMPT), HumanMessage(content=text)]
    response = llm.invoke(messages)
    return parse_json_response(response.content)


def score_fit(profile: dict, job_desc: str) -> dict:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    messages = [
        SystemMessage(content=FIT_PROMPT),
        HumanMessage(content=f"Candidate profile:\n{json.dumps(profile, indent=2)}\n\nJob description:\n{job_desc}"),
    ]
    response = llm.invoke(messages)
    return parse_json_response(response.content)


def format_report_markdown(profile: dict, fit: dict = None) -> str:
    skills_lang = ", ".join(profile.get("skills", {}).get("languages", []))
    md = f"""# Resume Analysis Report: {profile.get('name', 'Candidate')}

* **Current Title:** {profile.get('current_title', 'N/A')}
* **Years of Experience:** {profile.get('years_experience', 'N/A')}
* **Email:** {profile.get('email', 'N/A')}
* **Location:** {profile.get('location', 'N/A')}
* **Key Languages:** {skills_lang if skills_lang else 'N/A'}

### Summary
{profile.get('summary', 'N/A')}
"""

    if fit:
        fit_label = fit.get("fit_label", "N/A")
        label_emoji = {"Excellent": "🟢", "Good": "🟡", "Fair": "🟠", "Poor": "🔴"}.get(fit_label, "⚪")
        md += f"""
---

## Job Fit Analysis

* **Fit Score:** {label_emoji} {fit.get('fit_score', 'N/A')}/100 ({fit_label})
* **Recommendation:** {fit.get('recommendation', 'N/A')}
* **Strengths:** {', '.join(fit.get('strengths', []))}
* **Gaps:** {', '.join(fit.get('gaps', ['None identified']))}

### Reason
{fit.get('recommendation_reason', 'N/A')}
"""
    return md


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "resume_analysis.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "resume_analysis.docx") -> str:
    doc = Document()
    doc.add_heading("Resume Analysis Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "resume_analysis.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resume Analysis Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


SAMPLE_RESUME = """
Jane Doe
jane.doe@email.com | +1 (555) 123-4567 | San Francisco, CA
[linkedin.com/in/janedoe](https://linkedin.com/in/janedoe) | [github.com/janedoe](https://github.com/janedoe)

SUMMARY
Senior Python developer with 7 years of experience building scalable web applications
and data pipelines. Led teams of 5-8 engineers at Series B startups.

EXPERIENCE
Senior Software Engineer | TechCorp Inc. | 2021-present
- Architected microservices platform handling 10M requests/day using FastAPI + Kubernetes
- Reduced API latency by 40% through Redis caching and async optimization
- Led migration from monolith to microservices (12-month project, 5 engineers)

Software Engineer | DataFlow Systems | 2018-2021
- Built ML data pipelines processing 500GB/day using Apache Spark and Airflow
- Developed REST APIs with Django REST Framework serving 50k daily users

SKILLS
Languages: Python, JavaScript, SQL, Bash
Frameworks: FastAPI, Django, React, Spark
Tools: Docker, Kubernetes, Redis, PostgreSQL, Git, Airflow
Cloud: AWS (EC2, S3, RDS, Lambda)

EDUCATION
B.S. Computer Science | UC Berkeley | 2017

CERTIFICATIONS
AWS Solutions Architect Associate
"""


def main():
    parser = argparse.ArgumentParser(description="Resume Parser Agent")
    parser.add_argument("--resume", help="Path to resume file (.txt or .pdf)")
    parser.add_argument("--job-desc", help="Job description to match against")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    if args.resume and os.path.exists(args.resume):
        text = read_resume_text(args.resume)
    else:
        text = task_prompt if task_prompt else SAMPLE_RESUME

    job_desc = args.job_desc or (task_prompt if args.resume else None)

    profile = parse_resume(text)
    fit = score_fit(profile, job_desc) if job_desc else None

    report_md = format_report_markdown(profile, fit)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or job_desc or "", report_md)


if __name__ == "__main__":
    main()
