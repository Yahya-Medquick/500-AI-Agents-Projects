"""
PDF Q&A Agent using LlamaIndex and Gemini Key Rotation.

Loads a PDF, indexes it, and answers questions about its content.
Supports CLI flags as well as environment execution via TASK_PROMPT.
"""

import argparse
import base64
import os
import random
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex
from llama_index.llms.gemini import Gemini
from pptx import Presentation

load_dotenv()


def get_gemini_llm_llama_index(model_name: str = "models/gemini-3.5-flash", temperature: float = 0.0) -> Gemini:
    """
    Returns a LlamaIndex Gemini LLM instance configured with an active key
    drawn from a comma-separated GEMINI_API_KEY environment secret.
    """
    raw_env = os.getenv("GEMINI_API_KEY", "")
    keys = [k.strip() for k in raw_env.split(",") if k.strip()]
    
    if not keys:
        raise ValueError("GEMINI_API_KEY is not set or empty.")

    selected_key = random.choice(keys)
    return Gemini(model_name=model_name, api_key=selected_key, temperature=temperature)


def build_index(pdf_path: str) -> VectorStoreIndex:
    print(f"📄 Loading and indexing {pdf_path}...")
    reader = SimpleDirectoryReader(input_files=[pdf_path])
    docs = reader.load_data()
    index = VectorStoreIndex.from_documents(docs)
    print(f"✅ Indexed {len(docs)} document chunk(s)")
    return index


def single_question(index: VectorStoreIndex, question: str) -> str:
    llm = get_gemini_llm_llama_index()
    query_engine = index.as_query_engine(llm=llm, similarity_top_k=5)
    response = query_engine.query(question)
    
    output = str(response.response)
    if hasattr(response, "source_nodes") and response.source_nodes:
        output += f"\n\n---\n**Sources**: {len(response.source_nodes)} document chunk(s) referenced."
    return output


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "qa_report.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "qa_report.docx") -> str:
    doc = Document()
    doc.add_heading("PDF Q&A Response", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "qa_report.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PDF Q&A Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="PDF Q&A Agent")
    parser.add_argument("--pdf", help="Path to PDF file")
    parser.add_argument("--question", help="Question about PDF content")
    parser.add_argument("--query", help="General execution task")
    args = parser.parse_args()

    pdf_path = args.pdf or os.getenv("PDF_PATH")
    question = args.question or args.query or os.getenv("TASK_PROMPT") or "Summarize the key contents of this document."

    if not pdf_path or not os.path.exists(pdf_path):
        print("\n---REPORT_START---")
        print("⚠️ **Error**: No valid `--pdf` file path was provided or the file was not found.")
        print("---REPORT_END---\n")
        return

    index = build_index(pdf_path)
    answer = single_question(index, question)

    print("\n---REPORT_START---")
    print(answer)
    print("---REPORT_END---\n")

    handle_on_demand_export(question, answer)


if __name__ == "__main__":
    main()
