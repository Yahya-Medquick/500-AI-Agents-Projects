"""
Meeting Notes Agent using Gemini Key Rotation.

Converts meeting transcript text into structured meeting notes:
summary, action items, decisions, and follow-ups.
"""

import argparse
import base64
import json
import os
import re
from datetime import date, datetime

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

NOTES_PROMPT = """You are a professional meeting note-taker. Convert the meeting transcript into structured notes as JSON:
{
  "meeting_title": "inferred title",
  "date": "today or mentioned date",
  "participants": ["name1", "name2"],
  "duration_estimate": "X minutes",
  "summary": "2-3 sentence executive summary",
  "key_decisions": ["decision 1", "decision 2"],
  "action_items": [
    {"task": "description", "owner": "person name or TBD", "due": "date or timeframe or TBD"}
  ],
  "discussion_topics": ["topic 1", "topic 2"],
  "blockers": ["blocker 1 or none"],
  "next_meeting": "scheduled time or TBD",
  "follow_up_questions": ["question needing resolution"]
}
Return only valid JSON, no markdown formatting."""

SAMPLE_TRANSCRIPT = """
Sarah: Alright everyone, let's get started. It's Monday the 3rd and we have John, Mike, and Lisa here.

John: Thanks Sarah. So the main thing I wanted to cover is the Q4 product roadmap.
We need to decide on the feature freeze date.

Sarah: I think we should freeze by November 15th. That gives QA three weeks before the holiday release.

Mike: That works for me. But we still need to finalize the payment integration. Lisa, where are you on that?

Lisa: I'm about 60% done. I need the API docs from the payment provider. I've emailed them twice but haven't heard back.

John: I'll escalate that today. I'll reach out to our account manager at PaymentCo. That's blocking us.

Sarah: Okay, so John will handle the PaymentCo escalation by end of today. Lisa continues on payment integration, targeting completion by November 10th.

Mike: I can help with testing once Lisa has a draft ready. Let's say I start testing November 11th.

Sarah: Great. Also, we decided to cut the social login feature from this release. Too risky to add now.

John: Agreed. We'll put it in Q1 backlog.

Sarah: Any other blockers? No? Okay. Same time next week, November 10th.
"""


def parse_json_response(text: str) -> dict:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    return json.loads(cleaned)


def generate_meeting_notes(transcript: str) -> dict:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0)
    messages = [
        SystemMessage(content=NOTES_PROMPT),
        HumanMessage(content=f"Meeting transcript:\n\n{transcript}"),
    ]
    response = llm.invoke(messages)
    return parse_json_response(response.content)


def format_notes(notes: dict) -> str:
    lines = [
        f"# {notes.get('meeting_title', 'Meeting Notes')}",
        f"**Date:** {notes.get('date', date.today().isoformat())}  |  **Duration:** {notes.get('duration_estimate', 'N/A')}",
        f"**Participants:** {', '.join(notes.get('participants', []))}",
        "",
        "## Summary",
        notes.get("summary", ""),
        "",
        "## Key Decisions",
        *[f"- {d}" for d in notes.get("key_decisions", [])],
        "",
        "## Action Items",
    ]
    for item in notes.get("action_items", []):
        lines.append(f"- [ ] **{item.get('task', 'Task')}** — Owner: {item.get('owner', 'TBD')} | Due: {item.get('due', 'TBD')}")

    if notes.get("blockers"):
        lines += ["", "## Blockers", *[f"- {b}" for b in notes["blockers"]]]

    if notes.get("next_meeting") and notes["next_meeting"] != "TBD":
        lines += ["", f"**Next Meeting:** {notes['next_meeting']}"]

    return "\n".join(lines)


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "meeting_notes.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "meeting_notes.docx") -> str:
    doc = Document()
    doc.add_heading("Meeting Notes Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "meeting_notes.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Meeting Notes Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Meeting Notes Agent")
    group = parser.add_mutually_exclusive_group(required=False)
    group.add_argument("--transcript", help="Path to transcript text file")
    group.add_argument("--text", help="Transcript text directly")
    parser.add_argument("--query", help="General execution prompt/task")
    parser.add_argument("--output", default="meeting_notes.md", help="Markdown output path")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    if args.transcript and os.path.exists(args.transcript):
        with open(args.transcript, "r", encoding="utf-8", errors="ignore") as f:
            transcript = f.read()
    elif args.text:
        transcript = args.text
    elif task_prompt:
        transcript = task_prompt
    else:
        transcript = SAMPLE_TRANSCRIPT

    notes = generate_meeting_notes(transcript)
    formatted = format_notes(notes)

    print("\n---REPORT_START---")
    print(formatted)
    print("---REPORT_END---\n")

    # Save to file
    output_file = args.output
    if os.path.exists(output_file) and args.output == "meeting_notes.md":
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = f"meeting_notes_{stamp}.md"
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(formatted)

    handle_on_demand_export(task_prompt, formatted)


if __name__ == "__main__":
    main()
