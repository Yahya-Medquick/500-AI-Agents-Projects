"""
Travel Planner Agent using CrewAI and Gemini Key Rotation.

Multi-agent crew that creates personalized travel itineraries:
- Destination Researcher: gathers destination info
- Activity Planner: creates day-by-day activities
- Budget Analyst: estimates costs
"""

import argparse
import base64
import os
import re

from crewai import Agent, Crew, Process, Task
from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()


def build_travel_crew(destination: str, days: int, budget: float, interests: str) -> str:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.4)

    researcher = Agent(
        role="Destination Researcher",
        goal=f"Research {destination} and provide key travel insights",
        backstory="Expert travel journalist who has visited 100+ countries. Knows the best hidden gems and practical tips.",
        llm=llm,
        verbose=False,
    )

    planner = Agent(
        role="Travel Itinerary Planner",
        goal=f"Create a detailed {days}-day itinerary for {destination}",
        backstory="Luxury travel consultant with 15 years of experience crafting personalized itineraries.",
        llm=llm,
        verbose=False,
    )

    budget_analyst = Agent(
        role="Travel Budget Analyst",
        goal=f"Estimate realistic costs for the trip within ${budget} budget",
        backstory="Financial travel advisor who helps travelers maximize experiences within budget.",
        llm=llm,
        verbose=False,
    )

    research_task = Task(
        description=f"""Research {destination} for a {days}-day trip.
Cover: best time to visit, neighborhoods to stay in, must-see attractions,
local food scene, transportation tips, and cultural customs to know.
Traveler interests: {interests}""",
        agent=researcher,
        expected_output="Destination brief with key areas, attractions, food, and practical tips",
    )

    planning_task = Task(
        description=f"""Create a {days}-day itinerary for {destination}.
Budget: ${budget} total. Interests: {interests}.
Include morning/afternoon/evening activities, specific restaurant recommendations,
and travel time between locations. Make it achievable and enjoyable.""",
        agent=planner,
        expected_output=f"Day-by-day {days}-day itinerary with activities, meals, and timing",
        context=[research_task],
    )

    budget_task = Task(
        description=f"""Provide a budget breakdown for the {days}-day {destination} trip.
Total budget: ${budget}. Include: flights (estimate), accommodation, food, activities,
transportation. Flag if budget is tight and suggest adjustments.""",
        agent=budget_analyst,
        expected_output="Itemized budget breakdown with daily averages and money-saving tips",
        context=[research_task, planning_task],
    )

    crew = Crew(
        agents=[researcher, planner, budget_analyst],
        tasks=[research_task, planning_task, budget_task],
        process=Process.sequential,
        verbose=False,
    )

    return str(crew.kickoff())


def format_report_markdown(destination: str, days: int, budget: float, itinerary_text: str) -> str:
    return f"""# Travel Itinerary: {destination} ({days} Days)
**Budget:** ${budget} USD

---

{itinerary_text}
"""


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "travel_itinerary.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "travel_itinerary.docx") -> str:
    doc = Document()
    doc.add_heading("Travel Itinerary Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "travel_itinerary.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Travel Itinerary Summary"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Travel Planner Agent")
    parser.add_argument("--destination", help="Travel destination")
    parser.add_argument("--days", type=int, help="Number of days")
    parser.add_argument("--budget", type=float, help="Total budget in USD")
    parser.add_argument("--interests", help="Traveler interests")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    destination = args.destination or "Tokyo, Japan"
    days = args.days or 7
    budget = args.budget or 3000.0
    interests = args.interests or "food, culture, history"

    if task_prompt and not args.destination:
        destination = task_prompt

    itinerary = build_travel_crew(destination, days, budget, interests)
    report_md = format_report_markdown(destination, days, budget, itinerary)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or destination, report_md)


if __name__ == "__main__":
    main()
