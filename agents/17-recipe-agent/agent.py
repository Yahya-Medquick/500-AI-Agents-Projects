"""
Recipe Recommendation Agent using Agno-style single agent and Gemini Key Rotation.

Suggests recipes based on available ingredients, dietary restrictions,
and time constraints.
"""

import argparse
import base64
import json
import os
import re

from docx import Document
from dotenv import load_dotenv
from fpdf import FPDF
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

# Import the key rotator helper at the top
from scripts.gemini_rotator import get_gemini_llm

load_dotenv()

RECIPE_PROMPT = """You are a professional chef and nutritionist. Given available ingredients and constraints,
suggest 3 recipes. Return JSON:
{
  "recipes": [
    {
      "name": "Recipe Name",
      "cuisine": "Italian/Asian/etc",
      "difficulty": "Easy/Medium/Hard",
      "prep_time": "X minutes",
      "cook_time": "X minutes",
      "servings": N,
      "ingredients_needed": ["ingredient (amount)"],
      "missing_ingredients": ["optional additions"],
      "instructions": ["Step 1: ...", "Step 2: ..."],
      "nutrition_per_serving": {"calories": N, "protein": "Xg", "carbs": "Xg", "fat": "Xg"},
      "tips": "Chef's tip"
    }
  ],
  "recommended": "Recipe Name (best match)"
}
Return only valid JSON."""


def parse_json_response(text: str) -> dict:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)
    return json.loads(cleaned)


def get_recipes(ingredients: list[str], diet: str, time_limit: int, servings: int) -> dict:
    # Uses key rotator with gemini-3.5-flash
    llm = get_gemini_llm(model="gemini-3.5-flash", temperature=0.5)

    constraints = []
    if diet:
        constraints.append(f"Dietary restriction: {diet}")
    if time_limit:
        constraints.append(f"Max total time: {time_limit} minutes")
    if servings:
        constraints.append(f"Servings needed: {servings}")

    messages = [
        SystemMessage(content=RECIPE_PROMPT),
        HumanMessage(content=f"Available ingredients: {', '.join(ingredients)}\n\nConstraints:\n{chr(10).join(constraints) if constraints else 'None'}"),
    ]

    response = llm.invoke(messages)
    return parse_json_response(response.content)


def format_report_markdown(recipe_data: dict, ingredients: list[str]) -> str:
    recipes = recipe_data.get("recipes", [])
    recommended = recipe_data.get("recommended", "N/A")

    md_content = f"# Recipe Recommendations Report\n\n"
    md_content += f"**Available Ingredients:** {', '.join(ingredients)}\n"
    md_content += f"**Recommended Option:** {recommended}\n\n---\n\n"

    for r in recipes:
        md_content += f"## 🍽️ {r.get('name', 'Recipe')} ({r.get('cuisine', 'N/A')})\n"
        md_content += f"* **Prep Time:** {r.get('prep_time', 'N/A')} | **Cook Time:** {r.get('cook_time', 'N/A')} | **Difficulty:** {r.get('difficulty', 'N/A')}\n"
        md_content += f"* **Servings:** {r.get('servings', 'N/A')}\n\n"
        
        md_content += "### 📝 Ingredients\n"
        for ing in r.get("ingredients_needed", []):
            md_content += f"* {ing}\n"
        if r.get("missing_ingredients"):
            md_content += f"\n*Optional Additions:* {', '.join(r['missing_ingredients'])}\n"

        md_content += "\n### 👨‍🍳 Instructions\n"
        for i, step in enumerate(r.get("instructions", []), 1):
            md_content += f"{i}. {step}\n"

        n = r.get("nutrition_per_serving", {})
        if n:
            md_content += f"\n**Nutrition (per serving):** {n.get('calories', '?')} cal | Protein: {n.get('protein', '?')} | Carbs: {n.get('carbs', '?')} | Fat: {n.get('fat', '?')}\n"

        if r.get("tips"):
            md_content += f"\n💡 *Chef's Tip:* {r['tips']}\n"
        
        md_content += "\n---\n\n"

    return md_content


# --- ON-DEMAND EXPORT FUNCTIONS ---

def export_pdf(text: str, filename: str = "recipe_recommendations.pdf") -> str:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    clean_text = text.encode('latin-1', 'replace').decode('latin-1')
    pdf.write(5, clean_text)
    pdf.output(filename)
    return filename


def export_docx(text: str, filename: str = "recipe_recommendations.docx") -> str:
    doc = Document()
    doc.add_heading("Recipe Recommendations Report", level=1)
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(filename)
    return filename


def export_pptx(text: str, filename: str = "recipe_recommendations.pptx") -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recipe Recommendations"
    slide.placeholders[1].text = text[:1200]
    prs.save(filename)
    return filename


def handle_on_demand_export(prompt: str, report_text: str):
    prompt_lower = prompt.lower()
    generated_file = None

    if any(k in prompt_lower for k in ["docx", "word", "doc"]):
        generated_file = export_docx(report_text)
    elif "pdf" in prompt_lower:
        generated_file = export_pdf(report_text)
    elif any(k in prompt_lower for k in ["pptx", "presentation", "slides", "powerpoint"]):
        generated_file = export_pptx(report_text)

    if generated_file and os.path.exists(generated_file):
        with open(generated_file, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            print(f"\n---FILE_EXPORT_START:{generated_file}---")
            print(b64_str)
            print("---FILE_EXPORT_END---\n")


def main():
    parser = argparse.ArgumentParser(description="Recipe Recommendation Agent")
    parser.add_argument("--ingredients", help="Comma-separated ingredients")
    parser.add_argument("--diet", default="", help="Dietary restriction (vegan, vegetarian, gluten-free, keto, etc.)")
    parser.add_argument("--time", type=int, default=0, help="Max cooking time in minutes")
    parser.add_argument("--servings", type=int, default=2, help="Number of servings")
    parser.add_argument("--query", help="General execution prompt/task")
    args = parser.parse_args()

    task_prompt = args.query or os.getenv("TASK_PROMPT") or ""

    if args.ingredients:
        ingredients_str = args.ingredients
    elif task_prompt:
        ingredients_str = task_prompt
    else:
        ingredients_str = "chicken breast, garlic, lemon, olive oil, rosemary, potatoes"

    ingredients = [i.strip() for i in ingredients_str.split(",")]

    result = get_recipes(ingredients, args.diet, args.time, args.servings)
    report_md = format_report_markdown(result, ingredients)

    print("\n---REPORT_START---")
    print(report_md)
    print("---REPORT_END---\n")

    handle_on_demand_export(task_prompt or ingredients_str, report_md)


if __name__ == "__main__":
    main()
